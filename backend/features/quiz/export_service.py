"""
Export Service — generate PDF, DOCX, PPTX, XLSX reports for quiz sessions.
"""
from __future__ import annotations

import io
import os
import math
from dataclasses import dataclass, field
from datetime import datetime
from typing import Any, Dict, List, Optional, Tuple

from fastapi import HTTPException
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy.orm import joinedload, selectinload

from persistence.models.quiz import (
    Quiz, QuizSession, QuizType, Question, QuestionType, Answer
)
from features.quiz.answer_service_async import AnswerServiceAsync

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------

_LOGO_PATH = os.path.join(
    os.path.dirname(__file__),          # features/quiz/
    "..", "..",                         # backend/
    "uploads", "logo.png"
)
# Fallback: frontend asset (symlinked or copied at startup is not guaranteed,
# so we try the frontend source tree too)
_LOGO_FALLBACK = os.path.join(
    os.path.dirname(__file__),
    "..", "..", "..",
    "frontend", "src", "assets", "logo.png"
)


def _logo_path() -> Optional[str]:
    for p in (_LOGO_PATH, _LOGO_FALLBACK):
        resolved = os.path.normpath(p)
        if os.path.exists(resolved):
            return resolved
    return None


_UPLOADS_BASE = os.path.normpath(os.path.join(
    os.path.dirname(__file__), "..", "..", "uploads"
))

def _resolve_upload_path(stored: Optional[str]) -> Optional[str]:
    """Convert a stored upload path/URL to an absolute local filesystem path.

    Handles several formats:
    - Bare relative path:  "2/13/file.png"           → uploads/images/2/13/file.png
    - With images prefix:  "images/2/13/file.png"    → uploads/images/2/13/file.png
    - Full upload path:    "uploads/images/..."       → uploads/images/...
    - URL path:            "/api/uploads/images/..."  → uploads/images/...
    """
    if not stored:
        return None
    p = stored.lstrip("/")
    # Strip URL prefix
    if p.startswith("api/"):
        p = p[4:]
    # Now p may be: "uploads/images/...", "images/...", or bare "tenant/quiz/file.png"
    if p.startswith("uploads/"):
        full = os.path.normpath(os.path.join(os.path.dirname(__file__), "..", "..", p))
    elif p.startswith("images/"):
        full = os.path.normpath(os.path.join(os.path.dirname(__file__), "..", "..", "uploads", p))
    else:
        # Bare relative path — assume it lives under uploads/images/
        full = os.path.normpath(os.path.join(os.path.dirname(__file__), "..", "..", "uploads", "images", p))
    return full if os.path.exists(full) else None


# ---------------------------------------------------------------------------
# Data container
# ---------------------------------------------------------------------------

@dataclass
class QuestionExport:
    id: int
    text: str
    question_type: str
    points: int
    max_time_seconds: Optional[int]
    options: Optional[List[str]]
    correct_answer_index: Optional[int]
    answer_distribution: List[int]
    total_answers: int
    word_frequencies: Optional[Dict[str, int]] = None   # word cloud only
    text_answers: Optional[List[str]] = None             # single_line / paragraph only
    question_image_path: Optional[str] = None
    option_image_paths: Optional[Dict[str, str]] = None  # "A"/"B"/"C"/"D" -> local path


@dataclass
class LeaderboardEntryExport:
    rank: int
    display_name: str
    score: int
    time_taken_seconds: Optional[float]


@dataclass
class ExportData:
    session_id: int
    quiz_title: str
    quiz_type: str           # "quiz" | "poll"
    total_questions: int
    total_participants: int
    generated_at: datetime
    questions: List[QuestionExport] = field(default_factory=list)
    leaderboard: List[LeaderboardEntryExport] = field(default_factory=list)


# ---------------------------------------------------------------------------
# Colour palette helpers
# ---------------------------------------------------------------------------

_PALETTE_HEX = [
    "#1890ff", "#13c2c2", "#52c41a", "#faad14",
    "#f5222d", "#722ed1", "#eb2f96", "#fa8c16",
]
_CORRECT_HEX = "#52c41a"
_WRONG_HEX   = "#1890ff"


def _hex2rgb(h: str) -> Tuple[int, int, int]:
    h = h.lstrip("#")
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


# ---------------------------------------------------------------------------
# Word cloud image (wordcloud library → Pillow PNG)
# ---------------------------------------------------------------------------

def _make_wordcloud_png(word_freq: Dict[str, int], width: int = 560, height: int = 300) -> bytes:
    """Render a real word cloud PNG using the wordcloud library."""
    if not word_freq:
        return b""
    try:
        from wordcloud import WordCloud
        wc = WordCloud(
            width=width, height=height,
            background_color="white",
            colormap="tab10",
            max_words=80,
            prefer_horizontal=0.8,
        )
        wc.generate_from_frequencies(word_freq)
        buf = io.BytesIO()
        wc.to_image().save(buf, format="PNG")
        return buf.getvalue()
    except Exception:
        return b""


# ---------------------------------------------------------------------------
# Pillow PNG chart helpers  (bar + pie — used by DOCX)
# ---------------------------------------------------------------------------

def _make_bar_chart_png(
    question_text: str,
    options: List[str],
    distribution: List[int],
    correct_idx: Optional[int],
    width: int = 560,
) -> bytes:
    """Draw a horizontal bar chart PNG using Pillow."""
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        return b""

    row_h = 44
    margin_left = 160
    margin_right = 80
    margin_top = 10
    margin_bottom = 10
    height = margin_top + len(options) * row_h + margin_bottom

    img = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(img)

    total = sum(distribution) or 1
    bar_area = width - margin_left - margin_right

    for i, (opt, cnt) in enumerate(zip(options, distribution)):
        y = margin_top + i * row_h
        bar_w = int(bar_area * cnt / total)
        color = _CORRECT_HEX if i == correct_idx else _WRONG_HEX
        r, g, b = _hex2rgb(color)

        draw.rectangle([margin_left, y + 6, margin_left + max(bar_w, 2), y + row_h - 6],
                       fill=(r, g, b))

        label = (opt[:18] + "…") if len(opt) > 20 else opt
        draw.text((4, y + row_h // 2 - 7), label, fill=(0, 0, 0))

        pct = cnt / total * 100
        draw.text((margin_left + bar_w + 4, y + row_h // 2 - 7),
                  f"{cnt} ({pct:.0f}%)", fill=(80, 80, 80))

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _make_pie_chart_png(
    labels: List[str],
    values: List[int],
    title: str = "",
    size: int = 300,
) -> bytes:
    """Draw a simple pie chart PNG using Pillow."""
    try:
        from PIL import Image, ImageDraw
    except ImportError:
        return b""

    total = sum(values) or 1
    legend_h = 20 * len(labels) + 10
    img = Image.new("RGB", (size, size + legend_h), "white")
    draw = ImageDraw.Draw(img)

    cx, cy, r = size // 2, size // 2, size // 2 - 10
    start = -90.0
    colors = [_hex2rgb(c) for c in _PALETTE_HEX]
    for i, (lbl, val) in enumerate(zip(labels, values)):
        sweep = val / total * 360
        draw.pieslice([cx - r, cy - r, cx + r, cy + r],
                      start=start, end=start + sweep,
                      fill=colors[i % len(colors)])
        start += sweep

    for i, (lbl, val) in enumerate(zip(labels, values)):
        ly = size + 5 + i * 20
        r2, g2, b2 = colors[i % len(colors)]
        draw.rectangle([5, ly + 3, 17, ly + 15], fill=(r2, g2, b2))
        draw.text((22, ly), f"{lbl} ({val / total * 100:.0f}%)", fill=(0, 0, 0))

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _make_wc_bar_png(word_freq: Dict[str, int], top: int = 15, width: int = 560) -> bytes:
    """Horizontal bar chart for word cloud top words."""
    items = sorted(word_freq.items(), key=lambda x: -x[1])[:top]
    if not items:
        return b""
    labels = [w for w, _ in items]
    counts = [c for _, c in items]
    return _make_bar_chart_png("", labels, counts, correct_idx=None, width=width)


# ---------------------------------------------------------------------------
# PDF builder (reportlab PLATYPUS + native charts + header/footer)
# ---------------------------------------------------------------------------

def _build_pdf(data: ExportData) -> bytes:
    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import cm
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
            HRFlowable, KeepTogether, PageBreak, Image as RLImage
        )
        from reportlab.graphics.shapes import Drawing, Rect, String
        from reportlab.graphics.charts.barcharts import HorizontalBarChart
        from reportlab.graphics.charts.piecharts import Pie
        from reportlab.graphics import renderPDF
        from reportlab.pdfgen import canvas as rl_canvas
    except ImportError:
        raise HTTPException(status_code=501, detail="Export libraries not available")

    PAGE_W, PAGE_H = A4
    MARGIN = 1.5 * cm
    HEADER_H = 1.2 * cm   # reserved space at top for header band
    FOOTER_H = 1.0 * cm   # reserved space at bottom for page number

    logo = _logo_path()

    # ---- Canvas with header / footer drawn on every page ----
    class HeaderFooterCanvas(rl_canvas.Canvas):
        _hf_link_counter = 0  # class-level counter so names stay unique across state restores

        def __init__(self, *args, **kwargs):
            super().__init__(*args, **kwargs)
            self._saved_page_states: list = []

        def showPage(self):
            self._saved_page_states.append(dict(self.__dict__))
            self._startPage()

        def save(self):
            total_pages = len(self._saved_page_states)
            HeaderFooterCanvas._hf_link_counter = 0
            for state in self._saved_page_states:
                self.__dict__.update(state)
                self._draw_header_footer(total_pages)
                super().showPage()
            super().save()

        def _unique_link(self, url: str, rect) -> None:
            """Add a URL annotation using a globally unique name."""
            HeaderFooterCanvas._hf_link_counter += 1
            name = f"HFLink{HeaderFooterCanvas._hf_link_counter}"
            self._addAnnotation(
                self._makeLink(url, rect),
                name=name,
            )

        def _makeLink(self, url: str, rect):
            from reportlab.pdfbase.pdfdoc import PDFDictionary, PDFName, PDFArray, PDFString
            ann = PDFDictionary()
            ann["Type"] = PDFName("Annot")
            ann["Subtype"] = PDFName("Link")
            ann["Rect"] = PDFArray(self._absRect(rect, 0))
            A = PDFDictionary()
            A["Type"] = PDFName("Action")
            A["S"] = PDFName("URI")
            A["URI"] = PDFString(url)
            ann["A"] = A
            from reportlab.pdfbase.pdfdoc import PDFArray as _PA
            ann["Border"] = _PA([0, 0, 0])
            return ann

        def _draw_header_footer(self, total_pages: int):
            self.saveState()

            # ── Header line ───────────────────────────────────────────────
            header_y = PAGE_H - MARGIN - HEADER_H

            # logo (if available)
            logo_drawn_w = 0.0
            if logo:
                try:
                    logo_h_pt = HEADER_H - 4
                    logo_w_pt = logo_h_pt  # square assumption; reportlab scales it
                    self.drawImage(
                        logo,
                        MARGIN + 4, header_y + 2,
                        width=logo_w_pt, height=logo_h_pt,
                        preserveAspectRatio=True, mask="auto",
                    )
                    # hyperlink over logo image
                    self._unique_link(
                        'https://www.swaya.me',
                        (MARGIN + 4, header_y + 2,
                         MARGIN + 4 + logo_w_pt, header_y + 2 + logo_h_pt),
                    )
                    logo_drawn_w = logo_w_pt + 6
                except Exception:
                    pass

            # "Swaya.me (www.swaya.me)" text in header — brand blue
            self.setFillColorRGB(0.094, 0.565, 1.0)   # #1890ff
            self.setFont("Helvetica-Bold", 10)
            text_x = MARGIN + logo_drawn_w + 4
            text_y = header_y + HEADER_H * 0.3
            brand_text = "Swaya.me  (www.swaya.me)"
            self.drawString(text_x, text_y, brand_text)
            # hyperlink over brand text
            self._unique_link(
                'https://www.swaya.me',
                (text_x, text_y - 2, text_x + 160, text_y + 10),
            )

            # quiz title (right-aligned, truncated) — grey
            max_title = 55
            title_str = data.quiz_title if len(data.quiz_title) <= max_title else data.quiz_title[:max_title] + "…"
            self.setFillColorRGB(0.4, 0.4, 0.4)
            self.setFont("Helvetica", 9)
            self.drawRightString(PAGE_W - MARGIN - 4, header_y + HEADER_H * 0.3, title_str)

            # thin blue line below header content
            self.setStrokeColorRGB(0.094, 0.565, 1.0)
            self.setLineWidth(0.8)
            self.line(MARGIN, header_y, PAGE_W - MARGIN, header_y)

            # ── Footer line ────────────────────────────────────────────────
            footer_y = MARGIN - 2
            # thin grey line above footer content
            self.setStrokeColorRGB(0.75, 0.75, 0.75)
            self.setLineWidth(0.5)
            self.line(MARGIN, footer_y + FOOTER_H - 2, PAGE_W - MARGIN, footer_y + FOOTER_H - 2)

            self.setFillColorRGB(0.4, 0.4, 0.4)
            self.setFont("Helvetica", 8)
            page_num = self._saved_page_states.index(
                next(s for s in self._saved_page_states if s.get("_pageNumber") == self._pageNumber),
                0
            ) + 1 if hasattr(self, "_pageNumber") else 1
            self.drawCentredString(
                PAGE_W / 2, footer_y + 3,
                f"Page {self._pageNumber} of {total_pages}"
            )
            self.drawString(MARGIN + 4, footer_y + 3, "Generated by Swaya.me")
            self.drawRightString(PAGE_W - MARGIN - 4, footer_y + 3,
                                 data.generated_at.strftime("%d %b %Y"))

            self.restoreState()

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        leftMargin=MARGIN, rightMargin=MARGIN,
        topMargin=MARGIN + HEADER_H + 0.3 * cm,
        bottomMargin=MARGIN + FOOTER_H + 0.2 * cm,
    )
    styles = getSampleStyleSheet()
    story = []

    title_style = ParagraphStyle("Title2", parent=styles["Title"], fontSize=24, spaceAfter=4,
                                 textColor=colors.HexColor("#1890ff"))
    h2 = ParagraphStyle("H2", parent=styles["Heading2"], fontSize=13, spaceBefore=14, spaceAfter=4,
                        textColor=colors.HexColor("#1890ff"))
    body = styles["Normal"]
    small = ParagraphStyle("small", parent=body, fontSize=9, textColor=colors.grey)
    caption = ParagraphStyle("caption", parent=body, fontSize=10, textColor=colors.HexColor("#555555"),
                             spaceBefore=2, spaceAfter=6)

    # ---- Page 1: summary ----
    story.append(Paragraph(data.quiz_title, title_style))
    story.append(Paragraph(
        f"Session #{data.session_id}  ·  {data.generated_at.strftime('%d %b %Y, %H:%M')}",
        caption
    ))
    story.append(HRFlowable(width="100%", color=colors.HexColor("#1890ff"), thickness=3,
                            spaceAfter=8))

    # KPI row — large bold values
    kpi_data = [
        [
            Paragraph(f"<font size=22><b>{data.total_participants}</b></font><br/><font size=10 color='#555555'>Participants</font>", body),
            Paragraph(f"<font size=22><b>{data.total_questions}</b></font><br/><font size=10 color='#555555'>Questions</font>", body),
            Paragraph(f"<font size=22><b>{data.quiz_type.capitalize()}</b></font><br/><font size=10 color='#555555'>Type</font>", body),
        ]
    ]
    kpi_tbl = Table(kpi_data, colWidths=["33%", "33%", "34%"])
    kpi_tbl.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (0, 0), colors.HexColor("#dbeeff")),
        ("BACKGROUND", (1, 0), (1, 0), colors.HexColor("#d9f7d9")),
        ("BACKGROUND", (2, 0), (2, 0), colors.HexColor("#fff3d6")),
        ("BOX", (0, 0), (-1, -1), 0.5, colors.HexColor("#cccccc")),
        ("INNERGRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#dddddd")),
        ("ALIGN", (0, 0), (-1, -1), "CENTER"),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("TOPPADDING", (0, 0), (-1, -1), 14),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 14),
    ]))
    story.append(kpi_tbl)
    story.append(Spacer(1, 0.5 * cm))

    # Overall correct/incorrect pie (quiz type only) — use Pillow PNG for crisper rendering
    if data.quiz_type == "quiz":
        total_correct = 0
        total_wrong = 0
        for q in data.questions:
            if q.options and q.correct_answer_index is not None:
                for i, cnt in enumerate(q.answer_distribution):
                    if i == q.correct_answer_index:
                        total_correct += cnt
                    else:
                        total_wrong += cnt

        if total_correct + total_wrong > 0:
            pie_png = _make_pie_chart_png(
                [f"Correct ({total_correct})", f"Wrong ({total_wrong})"],
                [total_correct, total_wrong],
                title="Overall Results",
                size=240,
            )
            if pie_png:
                story.append(Paragraph("<b>Overall Results</b>", h2))
                story.append(RLImage(io.BytesIO(pie_png), width=7 * cm, height=8 * cm))
                story.append(Spacer(1, 0.3 * cm))

    # Per-question overview bar chart — use Pillow PNG for crisper rendering
    if data.questions:
        q_labels = [f"Q{i+1}" for i in range(len(data.questions))]
        q_counts = [sum(q.answer_distribution) for q in data.questions]
        overview_png = _make_bar_chart_png(
            "", q_labels, q_counts, correct_idx=None, width=520
        )
        if overview_png:
            story.append(Paragraph("<b>Response count per question</b>", h2))
            story.append(RLImage(io.BytesIO(overview_png), width=14 * cm,
                                 height=max(3, len(data.questions) * 0.55) * cm))

    story.append(PageBreak())

    # ---- Per-question pages ----
    for idx, q in enumerate(data.questions):
        total = sum(q.answer_distribution) or 1
        q_items = []

        q_items.append(HRFlowable(width="100%", color=colors.HexColor("#1890ff"),
                                  thickness=1, spaceBefore=6, spaceAfter=4))
        type_badge = q.question_type.upper()
        q_items.append(Paragraph(
            f"<b>Q{idx+1}. {q.text}</b>  <font color='#888888' size=8>[{type_badge}]</font>", h2))

        # Question image
        if q.question_image_path:
            try:
                q_items.append(RLImage(q.question_image_path, width=10*cm, height=6*cm))
                q_items.append(Spacer(1, 0.2*cm))
            except Exception:
                pass

        if q.options:
            # Pillow-based bar chart for crisp rendering
            bar_png = _make_bar_chart_png(
                q.text, q.options, q.answer_distribution,
                correct_idx=q.correct_answer_index, width=520,
            )
            if bar_png:
                bar_h = max(2.5, len(q.options) * 1.1) * cm
                q_items.append(RLImage(io.BytesIO(bar_png), width=14 * cm, height=bar_h))
                q_items.append(Spacer(1, 0.2 * cm))

            tbl_data = [["Option", "Votes", "%"]]
            for i, (opt, cnt) in enumerate(zip(q.options, q.answer_distribution)):
                marker = " ✓" if i == q.correct_answer_index else ""
                pct = cnt / total * 100
                tbl_data.append([f"{chr(65+i)}. {opt}{marker}", str(cnt), f"{pct:.0f}%"])
            style_cmds = [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1890ff")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("BOX", (0, 0), (-1, -1), 0.4, colors.HexColor("#cccccc")),
                ("INNERGRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#e0e0e0")),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f5f5f5")]),
                ("TOPPADDING", (0, 0), (-1, -1), 5),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
            ]
            # Highlight the correct answer row in green
            if q.correct_answer_index is not None:
                cr = q.correct_answer_index + 1  # +1 for header row
                style_cmds += [
                    ("BACKGROUND", (0, cr), (-1, cr), colors.HexColor("#e8f5e9")),
                    ("TEXTCOLOR", (0, cr), (-1, cr), colors.HexColor("#2e7d32")),
                    ("FONTNAME", (0, cr), (-1, cr), "Helvetica-Bold"),
                ]
            tbl = Table(tbl_data, colWidths=["60%", "20%", "20%"])
            tbl.setStyle(TableStyle(style_cmds))
            q_items.append(tbl)

            # Option images (after the MCQ table)
            if q.option_image_paths:
                for opt_key in ("A", "B", "C", "D"):
                    img_path = q.option_image_paths.get(opt_key)
                    if img_path:
                        try:
                            q_items.append(Paragraph(f"<b>Option {opt_key} image:</b>", small))
                            q_items.append(RLImage(img_path, width=8*cm, height=5*cm))
                            q_items.append(Spacer(1, 0.1*cm))
                        except Exception:
                            pass

            story.append(KeepTogether(q_items))

        elif q.question_type == 'word_cloud':
            # Add question title directly to story first
            story.extend(q_items)
            q_items = []  # clear so KeepTogether below is skipped

            if q.word_frequencies:
                # Real word cloud image — added directly (not in KeepTogether)
                wc_png = _make_wordcloud_png(q.word_frequencies, width=500, height=260)
                if wc_png:
                    story.append(Paragraph("<b>Word Cloud</b>", body))
                    story.append(RLImage(io.BytesIO(wc_png), width=14 * cm, height=7.3 * cm))
                    story.append(Spacer(1, 0.3 * cm))

                # Top-words bar chart added directly (not in KeepTogether)
                top_words = sorted(q.word_frequencies.items(), key=lambda x: -x[1])[:15]
                w_labels = [w for w, _ in top_words]
                w_counts = [c for _, c in top_words]
                max_wv = max(w_counts) if w_counts else 1
                dw3, dh3 = 420, max(60, len(w_labels) * 24 + 20)
                d4 = Drawing(dw3, dh3)
                bc3 = HorizontalBarChart()
                bc3.x = 60; bc3.y = 5; bc3.width = dw3 - 80; bc3.height = dh3 - 10
                bc3.data = [w_counts]
                bc3.categoryAxis.categoryNames = w_labels
                bc3.valueAxis.valueMin = 0
                bc3.valueAxis.valueMax = max_wv + 1
                bc3.bars[0].fillColor = colors.HexColor("#722ed1")
                d4.add(bc3)
                story.append(Paragraph("<b>Top word frequencies</b>", body))
                story.append(d4)
            else:
                story.append(Paragraph("<i>No responses submitted yet.</i>", body))

        else:
            if q.text_answers:
                q_items.append(Paragraph(
                    f"<i>Text responses ({len(q.text_answers)} received)</i>", body
                ))
                story.append(KeepTogether(q_items))
                for ans in q.text_answers:
                    story.append(Paragraph(f"• {ans}", body))
            else:
                q_items.append(Paragraph(
                    f"<i>Text response question — {sum(q.answer_distribution)} responses received</i>",
                    body
                ))
                story.append(KeepTogether(q_items))

        story.append(Spacer(1, 0.6 * cm))

    # ---- Leaderboard page ----
    if data.quiz_type == "quiz" and data.leaderboard:
        story.append(PageBreak())
        story.append(Paragraph("<b>Leaderboard</b>", title_style))
        story.append(Spacer(1, 0.3 * cm))

        top3 = data.leaderboard[:3]
        if len(top3) >= 1:
            podium_d = Drawing(400, 120)
            heights = [80, 60, 40]
            positions = [120, 20, 220]
            podium_colors = [
                colors.HexColor("#FFD700"),
                colors.HexColor("#C0C0C0"),
                colors.HexColor("#CD7F32"),
            ]
            for rank_0 in range(min(3, len(top3))):
                entry = top3[rank_0]
                h = heights[rank_0]
                x = positions[rank_0]
                c = podium_colors[rank_0]
                podium_d.add(Rect(x, 10, 60, h, fillColor=c, strokeColor=colors.white))
                podium_d.add(String(x + 30, 10 + h + 4, entry.display_name[:10],
                                    textAnchor="middle", fontSize=8))
                podium_d.add(String(x + 30, 10 + h // 2,
                                    f"#{entry.rank} ({entry.score})",
                                    textAnchor="middle", fontSize=9, fillColor=colors.black))
            story.append(podium_d)

        lb_data = [["Rank", "Name", "Score", "Time"]]
        rank_colors_map = {1: "#FFD700", 2: "#C0C0C0", 3: "#CD7F32"}
        for e in data.leaderboard:
            time_str = f"{e.time_taken_seconds:.1f}s" if e.time_taken_seconds is not None else "—"
            lb_data.append([str(e.rank), e.display_name, str(e.score), time_str])
        lb_tbl = Table(lb_data, colWidths=["10%", "50%", "20%", "20%"])
        style_cmds = [
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1890ff")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("BOX", (0, 0), (-1, -1), 0.4, colors.HexColor("#cccccc")),
            ("INNERGRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#e0e0e0")),
            ("TOPPADDING", (0, 0), (-1, -1), 5),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
            ("ALIGN", (0, 0), (0, -1), "CENTER"),
            ("ALIGN", (2, 0), (-1, -1), "RIGHT"),
        ]
        # Alternating rows only for entries beyond top-3 (guard against short leaderboards)
        if len(data.leaderboard) > 3:
            style_cmds.append(("ROWBACKGROUNDS", (0, 4), (-1, -1), [colors.white, colors.HexColor("#f5faff")]))
        for row_i, e in enumerate(data.leaderboard, start=1):
            hex_c = rank_colors_map.get(e.rank)
            if hex_c:
                style_cmds.append(("BACKGROUND", (0, row_i), (-1, row_i), colors.HexColor(hex_c)))
                style_cmds.append(("FONTNAME", (0, row_i), (-1, row_i), "Helvetica-Bold"))
        lb_tbl.setStyle(TableStyle(style_cmds))
        story.append(lb_tbl)

    doc.build(story, canvasmaker=HeaderFooterCanvas)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# DOCX builder (python-docx + Pillow PNGs)
# ---------------------------------------------------------------------------

def _build_docx(data: ExportData) -> bytes:
    try:
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor, Cm
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
    except ImportError:
        raise HTTPException(status_code=501, detail="Export libraries not available")

    doc = Document()

    # ---- Header (logo + title + divider) ----
    section = doc.sections[0]
    header = section.header
    header.is_linked_to_previous = False
    htbl = header.add_table(1, 2, Inches(6.5))
    htbl.style = "Table Grid"
    htbl.autofit = False
    # Remove border styling
    for row in htbl.rows:
        for cell in row.cells:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            tcBorders = OxmlElement("w:tcBorders")
            for side in ("top", "left", "bottom", "right"):
                border = OxmlElement(f"w:{side}")
                border.set(qn("w:val"), "none")
                tcBorders.append(border)
            tcPr.append(tcBorders)

    logo = _logo_path()
    left_cell = htbl.cell(0, 0)
    left_cell.width = Inches(0.6)
    lp = left_cell.paragraphs[0]
    lp.alignment = WD_ALIGN_PARAGRAPH.LEFT
    if logo:
        try:
            run = lp.add_run()
            run.add_picture(logo, width=Inches(0.45))
        except Exception:
            lp.add_run("Swaya.me")
    else:
        lp.add_run("Swaya.me").bold = True

    right_cell = htbl.cell(0, 1)
    rp = right_cell.paragraphs[0]
    rp.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    # Build a hyperlink run pointing to https://www.swaya.me
    r_id = header.part.relate_to(
        'https://www.swaya.me',
        'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink',
        is_external=True,
    )
    hyperlink_elem = OxmlElement('w:hyperlink')
    hyperlink_elem.set(qn('r:id'), r_id)
    run_elem = OxmlElement('w:r')
    rPr = OxmlElement('w:rPr')
    b_elem = OxmlElement('w:b')
    rPr.append(b_elem)
    color_elem = OxmlElement('w:color')
    color_elem.set(qn('w:val'), '1890FF')
    rPr.append(color_elem)
    sz_elem = OxmlElement('w:sz')
    sz_elem.set(qn('w:val'), '24')
    rPr.append(sz_elem)
    run_elem.append(rPr)
    t_elem = OxmlElement('w:t')
    t_elem.text = 'Swaya.me  (www.swaya.me)'
    run_elem.append(t_elem)
    hyperlink_elem.append(run_elem)
    rp._p.append(hyperlink_elem)

    # ---- Footer (page numbers) ----
    footer = section.footer
    footer.is_linked_to_previous = False
    fp = footer.paragraphs[0]
    fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    fp.add_run("Page ").font.size = Pt(8)
    # Insert PAGE field
    fldChar1 = OxmlElement("w:fldChar")
    fldChar1.set(qn("w:fldCharType"), "begin")
    instrText = OxmlElement("w:instrText")
    instrText.text = "PAGE"
    fldChar2 = OxmlElement("w:fldChar")
    fldChar2.set(qn("w:fldCharType"), "end")
    run_pg = fp.add_run()
    run_pg._r.append(fldChar1)
    run_pg._r.append(instrText)
    run_pg._r.append(fldChar2)
    run_pg.font.size = Pt(8)

    fp.add_run(" of ").font.size = Pt(8)

    fldChar3 = OxmlElement("w:fldChar")
    fldChar3.set(qn("w:fldCharType"), "begin")
    instrText2 = OxmlElement("w:instrText")
    instrText2.text = "NUMPAGES"
    fldChar4 = OxmlElement("w:fldChar")
    fldChar4.set(qn("w:fldCharType"), "end")
    run_np = fp.add_run()
    run_np._r.append(fldChar3)
    run_np._r.append(instrText2)
    run_np._r.append(fldChar4)
    run_np.font.size = Pt(8)

    fp.add_run(f"  ·  Generated by Swaya.me  ·  {data.generated_at.strftime('%d %b %Y')}").font.size = Pt(8)

    def _set_cell_bg(cell, hex_color: str):
        """Set background color of a table cell."""
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        shd = OxmlElement('w:shd')
        shd.set(qn('w:val'), 'clear')
        shd.set(qn('w:color'), 'auto')
        shd.set(qn('w:fill'), hex_color.lstrip('#'))
        tcPr.append(shd)

    def _add_colored_heading(doc, text: str, level: int, hex_color: str = "1890FF"):
        """Add a heading paragraph with custom color."""
        h = doc.add_heading(text, level)
        for run in h.runs:
            run.font.color.rgb = RGBColor(
                int(hex_color[0:2], 16),
                int(hex_color[2:4], 16),
                int(hex_color[4:6], 16),
            )
        return h

    # ---- Document body ----
    # Title block — blue background banner
    title_tbl = doc.add_table(rows=1, cols=1)
    title_tbl.style = "Table Grid"
    title_cell = title_tbl.cell(0, 0)
    _set_cell_bg(title_cell, "1890FF")
    title_para = title_cell.paragraphs[0]
    title_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
    title_run = title_para.add_run(data.quiz_title)
    title_run.bold = True
    title_run.font.size = Pt(18)
    title_run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    title_cell.add_paragraph()
    sub_run = title_cell.add_paragraph().add_run(
        f"Session #{data.session_id}  ·  {data.generated_at.strftime('%d %b %Y, %H:%M')}"
    )
    sub_run.font.size = Pt(10)
    sub_run.font.color.rgb = RGBColor(0xCC, 0xE5, 0xFF)
    doc.add_paragraph()

    # KPI summary table
    kpi_tbl = doc.add_table(rows=2, cols=3)
    kpi_tbl.style = "Table Grid"
    kpi_headers = ["Participants", "Questions", "Type"]
    kpi_values  = [str(data.total_participants), str(data.total_questions),
                   data.quiz_type.capitalize()]
    kpi_bg = ["DBEEFF", "D9F7D9", "FFF3D6"]
    for i, (h, v, bg) in enumerate(zip(kpi_headers, kpi_values, kpi_bg)):
        hcell = kpi_tbl.cell(0, i)
        hcell.text = h
        hcell.paragraphs[0].runs[0].bold = True
        _set_cell_bg(hcell, bg)
        vcell = kpi_tbl.cell(1, i)
        vp = vcell.paragraphs[0]
        vrun = vp.add_run(v)
        vrun.bold = True
        vrun.font.size = Pt(14)
        _set_cell_bg(vcell, bg)
    doc.add_paragraph()

    if data.quiz_type == "quiz":
        total_correct = sum(
            q.answer_distribution[q.correct_answer_index]
            for q in data.questions
            if q.options and q.correct_answer_index is not None
               and q.correct_answer_index < len(q.answer_distribution)
        )
        total_answers = sum(sum(q.answer_distribution) for q in data.questions if q.options)
        total_wrong = total_answers - total_correct
        if total_correct + total_wrong > 0:
            pie_png = _make_pie_chart_png(
                ["Correct", "Wrong"], [total_correct, total_wrong],
                title="Overall Correct vs Wrong"
            )
            if pie_png:
                _add_colored_heading(doc, "Overall Results", level=1)
                doc.add_picture(io.BytesIO(pie_png), width=Inches(3))

    _add_colored_heading(doc, "Question Results", level=1)
    for idx, q in enumerate(data.questions):
        _add_colored_heading(doc, f"Q{idx+1}. {q.text}", level=2)
        # Question image
        if q.question_image_path:
            try:
                doc.add_picture(q.question_image_path, width=Inches(4))
                doc.add_paragraph()
            except Exception:
                pass
        p = doc.add_paragraph()
        run = p.add_run(f"[{q.question_type.upper()}]  Total answers: {sum(q.answer_distribution)}")
        run.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
        run.font.size = Pt(9)

        if q.options:
            chart_png = _make_bar_chart_png(
                q.text, q.options, q.answer_distribution, q.correct_answer_index
            )
            if chart_png:
                doc.add_picture(io.BytesIO(chart_png), width=Inches(5.5))

            tbl2 = doc.add_table(rows=len(q.options) + 1, cols=3)
            tbl2.style = "Table Grid"
            # Header row — blue background
            for cell_text, cell in zip(["Option", "Votes", "%"], tbl2.rows[0].cells):
                cell.text = cell_text
                cell.paragraphs[0].runs[0].bold = True
                _set_cell_bg(cell, "1890FF")
                cell.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            total = sum(q.answer_distribution) or 1
            for i, (opt, cnt) in enumerate(zip(q.options, q.answer_distribution)):
                row = tbl2.rows[i + 1]
                marker = " ✓" if i == q.correct_answer_index else ""
                row.cells[0].text = f"{chr(65+i)}. {opt}{marker}"
                row.cells[1].text = str(cnt)
                row.cells[2].text = f"{cnt/total*100:.0f}%"
                # Highlight correct answer row green
                if i == q.correct_answer_index:
                    for cell in row.cells:
                        _set_cell_bg(cell, "E8F5E9")
                        for run in cell.paragraphs[0].runs:
                            run.bold = True
                            run.font.color.rgb = RGBColor(0x2E, 0x7D, 0x32)
                elif i % 2 == 1:
                    for cell in row.cells:
                        _set_cell_bg(cell, "F5F5F5")

            # Option images
            if q.option_image_paths:
                for opt_key in ("A", "B", "C", "D"):
                    img_path = q.option_image_paths.get(opt_key)
                    if img_path:
                        try:
                            doc.add_paragraph(f"Option {opt_key}:")
                            doc.add_picture(img_path, width=Inches(3))
                        except Exception:
                            pass

        elif q.question_type == 'word_cloud':
            if q.word_frequencies:
                wc_png = _make_wordcloud_png(q.word_frequencies, width=560, height=280)
                if wc_png:
                    doc.add_paragraph("Word Cloud:").bold = True
                    doc.add_picture(io.BytesIO(wc_png), width=Inches(5.5))
                wc_bar = _make_wc_bar_png(q.word_frequencies)
                if wc_bar:
                    doc.add_paragraph("Top word frequencies:").bold = True
                    doc.add_picture(io.BytesIO(wc_bar), width=Inches(5.5))
            else:
                doc.add_paragraph("No responses submitted yet.").italic = True
        else:
            doc.add_paragraph(
                f"Text response question — {sum(q.answer_distribution)} responses received"
            ).italic = True

        doc.add_paragraph()

    if data.quiz_type == "quiz" and data.leaderboard:
        doc.add_page_break()
        _add_colored_heading(doc, "Leaderboard", level=1)
        lb_tbl = doc.add_table(rows=len(data.leaderboard) + 1, cols=4)
        lb_tbl.style = "Table Grid"
        # Header row
        for cell_text, cell in zip(["Rank", "Name", "Score", "Time"], lb_tbl.rows[0].cells):
            cell.text = cell_text
            cell.paragraphs[0].runs[0].bold = True
            _set_cell_bg(cell, "1890FF")
            cell.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        rank_bg = {1: "FFD700", 2: "E8E8E8", 3: "CD7F32"}
        for i, e in enumerate(data.leaderboard):
            row = lb_tbl.rows[i + 1]
            time_str = f"{e.time_taken_seconds:.1f}s" if e.time_taken_seconds is not None else "—"
            row.cells[0].text = str(e.rank)
            row.cells[1].text = e.display_name
            row.cells[2].text = str(e.score)
            row.cells[3].text = time_str
            bg = rank_bg.get(e.rank, "F5FAFF" if i % 2 == 0 else "FFFFFF")
            for cell in row.cells:
                _set_cell_bg(cell, bg)
                if e.rank in rank_bg:
                    for run in cell.paragraphs[0].runs:
                        run.bold = True

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# PPTX builder (python-pptx native charts + header/footer)
# ---------------------------------------------------------------------------

def _build_pptx(data: ExportData) -> bytes:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.enum.text import PP_ALIGN
        from pptx.chart.data import ChartData
    except ImportError:
        raise HTTPException(status_code=501, detail="Export libraries not available")

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # blank

    W = prs.slide_width
    H = prs.slide_height

    HEADER_H = Inches(0.42)
    FOOTER_H = Inches(0.25)
    logo = _logo_path()

    def add_text_box(slide, text, left, top, width, height, bold=False, size=18, color=None):
        txb = slide.shapes.add_textbox(left, top, width, height)
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = RGBColor(*_hex2rgb(color))
        return txb

    def add_header_footer(slide, page_num: int, total_pages: int):
        """Draw header line + footer line on a slide (no filled rectangles)."""
        from pptx.enum.text import PP_ALIGN

        # ── Header ──
        # Logo in header
        logo_x = Inches(0.08)
        if logo:
            try:
                pil_logo = __import__("PIL.Image", fromlist=["Image"]).open(logo)
                aspect = pil_logo.width / pil_logo.height
                logo_h = HEADER_H - Inches(0.06)
                logo_w = int(logo_h * aspect)
                slide.shapes.add_picture(logo, logo_x, Inches(0.03), logo_w, logo_h)
                logo_x += logo_w + Inches(0.05)
            except Exception:
                pass

        # "Swaya.me" in header — brand blue (transparent background now)
        brand_tb = slide.shapes.add_textbox(logo_x, Inches(0.04), Inches(1.5), HEADER_H)
        brand_tf = brand_tb.text_frame
        brand_p = brand_tf.paragraphs[0]
        brand_r = brand_p.add_run()
        brand_r.text = "Swaya.me  (www.swaya.me)"
        brand_r.font.bold = True
        brand_r.font.size = Pt(10)
        brand_r.font.color.rgb = RGBColor(0x18, 0x90, 0xFF)
        brand_r.hyperlink.address = 'https://www.swaya.me'
        brand_tb.width = Inches(2.2)

        # Quiz title (right side of header)
        title_str = data.quiz_title[:60] + ("…" if len(data.quiz_title) > 60 else "")
        title_tb = slide.shapes.add_textbox(Inches(2), Inches(0.04), W - Inches(2.1), HEADER_H)
        title_tf = title_tb.text_frame
        title_p = title_tf.paragraphs[0]
        title_p.alignment = PP_ALIGN.RIGHT
        title_r = title_p.add_run()
        title_r.text = title_str
        title_r.font.size = Pt(9)
        title_r.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

        # ── Footer ──
        footer_tb = slide.shapes.add_textbox(0, H - FOOTER_H, W, FOOTER_H)
        footer_tf = footer_tb.text_frame
        footer_p = footer_tf.paragraphs[0]
        footer_p.alignment = PP_ALIGN.CENTER
        footer_r = footer_p.add_run()
        footer_r.text = f"Swaya.me  ·  Page {page_num} of {total_pages}  ·  {data.generated_at.strftime('%d %b %Y')}"
        footer_r.font.size = Pt(7)
        footer_r.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # Count total slides up front so we can pass total_pages
    # Slide count: 1 (summary) + len(questions) + (1 if leaderboard)
    total_slides = 1 + len(data.questions) + (1 if data.quiz_type == "quiz" and data.leaderboard else 0)

    # ---- Slide 1: Title / Summary ----
    slide1 = prs.slides.add_slide(blank_layout)
    add_header_footer(slide1, 1, total_slides)

    add_text_box(slide1, data.quiz_title,
                 Inches(0.5), HEADER_H + Inches(0.1), W - Inches(1), Inches(1.2),
                 bold=True, size=28)
    add_text_box(slide1,
                 f"Session #{data.session_id}  ·  {data.generated_at.strftime('%d %b %Y, %H:%M')}",
                 Inches(0.5), HEADER_H + Inches(1.2), W - Inches(1), Inches(0.5),
                 size=14, color="#888888")

    kpi_items = [
        (str(data.total_participants), "Participants", "DBEEFF", "1890FF"),
        (str(data.total_questions),    "Questions",    "D9F7D9", "52C41A"),
        (data.quiz_type.capitalize(),  "Type",         "FFF3D6", "FAAD14"),
    ]
    box_w = Inches(2.2)
    for i, (val, lbl, bg_hex, accent_hex) in enumerate(kpi_items):
        left = Inches(0.5) + i * (box_w + Inches(0.2))
        shp = slide1.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            left, HEADER_H + Inches(1.8), box_w, Inches(1.0)
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(
            int(bg_hex[0:2], 16), int(bg_hex[2:4], 16), int(bg_hex[4:6], 16)
        )
        shp.line.color.rgb = RGBColor(
            int(accent_hex[0:2], 16), int(accent_hex[2:4], 16), int(accent_hex[4:6], 16)
        )
        shp.line.width = Pt(1.5)
        tf = shp.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        r0 = p0.add_run()
        r0.text = val
        r0.font.bold = True
        r0.font.size = Pt(22)
        r0.font.color.rgb = RGBColor(
            int(accent_hex[0:2], 16), int(accent_hex[2:4], 16), int(accent_hex[4:6], 16)
        )
        p1 = tf.add_paragraph()
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = lbl
        r1.font.size = Pt(10)
        r1.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    if data.quiz_type == "quiz":
        total_correct = sum(
            q.answer_distribution[q.correct_answer_index]
            for q in data.questions
            if q.options and q.correct_answer_index is not None
               and q.correct_answer_index < len(q.answer_distribution)
        )
        total_answers = sum(sum(q.answer_distribution) for q in data.questions if q.options)
        total_wrong = total_answers - total_correct
        if total_correct + total_wrong > 0:
            cd = ChartData()
            cd.categories = ["Correct", "Wrong"]
            cd.add_series("Responses", [total_correct, total_wrong])
            chart = slide1.shapes.add_chart(
                XL_CHART_TYPE.PIE,
                Inches(0.5), HEADER_H + Inches(3.0), Inches(4), Inches(3),
                cd
            ).chart
            chart.series[0].points[0].format.fill.solid()
            chart.series[0].points[0].format.fill.fore_color.rgb = RGBColor(0x52, 0xc4, 0x1a)
            chart.series[0].points[1].format.fill.solid()
            chart.series[0].points[1].format.fill.fore_color.rgb = RGBColor(0xf5, 0x22, 0x2d)

    # ---- Per-question slides ----
    for idx, q in enumerate(data.questions):
        sl = prs.slides.add_slide(blank_layout)
        add_header_footer(sl, idx + 2, total_slides)
        add_text_box(sl, f"Q{idx+1}. {q.text}",
                     Inches(0.3), HEADER_H + Inches(0.05), W - Inches(0.6), Inches(1.0),
                     bold=True, size=16)

        content_top = HEADER_H + Inches(1.1)
        content_h = H - HEADER_H - FOOTER_H - Inches(1.2)

        # Question image — place to the right if present
        q_img_w = 0
        if q.question_image_path:
            try:
                from PIL import Image as PILImage
                pil = PILImage.open(q.question_image_path)
                aspect = pil.width / pil.height
                img_h = min(Inches(1.5), content_h)
                img_w = int(img_h * aspect)
                sl.shapes.add_picture(q.question_image_path, W - Inches(0.3) - img_w, HEADER_H + Inches(0.05), img_w, img_h)
                q_img_w = img_w + Inches(0.1)
            except Exception:
                pass

        if q.options and q.answer_distribution:
            cd2 = ChartData()
            cd2.categories = [
                (opt[:30] + "…") if len(opt) > 32 else opt
                for opt in q.options
            ]
            cd2.add_series("Votes", q.answer_distribution)
            chart2 = sl.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED,
                Inches(0.3), content_top, W - Inches(0.6), content_h,
                cd2
            ).chart
            if q.correct_answer_index is not None:
                for i in range(len(q.options)):
                    pt = chart2.series[0].points[i]
                    pt.format.fill.solid()
                    if i == q.correct_answer_index:
                        pt.format.fill.fore_color.rgb = RGBColor(0x52, 0xc4, 0x1a)
                    else:
                        pt.format.fill.fore_color.rgb = RGBColor(0x18, 0x90, 0xff)

            # Option images strip
            if q.option_image_paths:
                opt_keys = [k for k in ("A","B","C","D") if k in q.option_image_paths]
                if opt_keys:
                    strip_w = (W - Inches(0.6)) / len(opt_keys)
                    strip_top = H - FOOTER_H - Inches(1.2)
                    for i, ok in enumerate(opt_keys):
                        try:
                            sl.shapes.add_picture(
                                q.option_image_paths[ok],
                                Inches(0.3) + i * strip_w, strip_top,
                                strip_w - Inches(0.05), Inches(1.1)
                            )
                        except Exception:
                            pass

        elif q.question_type == 'word_cloud':
            if q.word_frequencies:
                # Word cloud image in top half, bar chart in bottom half
                wc_png = _make_wordcloud_png(q.word_frequencies, width=600, height=260)
                if wc_png:
                    sl.shapes.add_picture(
                        io.BytesIO(wc_png),
                        Inches(0.3), content_top,
                        W - Inches(0.6), content_h * 0.55
                    )
                    content_top += content_h * 0.55 + Inches(0.05)
                    content_h = content_h * 0.42

                top_wf = sorted(q.word_frequencies.items(), key=lambda x: -x[1])[:10]
                if top_wf:
                    cd3 = ChartData()
                    cd3.categories = [w for w, _ in top_wf]
                    cd3.add_series("Frequency", [c for _, c in top_wf])
                    sl.shapes.add_chart(
                        XL_CHART_TYPE.BAR_CLUSTERED,
                        Inches(0.3), content_top, W - Inches(0.6), content_h,
                        cd3
                    )
            else:
                add_text_box(sl,
                             "No responses submitted yet.",
                             Inches(0.3), content_top, W - Inches(0.6), Inches(1),
                             size=14, color="#888888")
        else:
            add_text_box(sl,
                         f"Text response question — {sum(q.answer_distribution)} responses received",
                         Inches(0.3), content_top, W - Inches(0.6), Inches(1),
                         size=14, color="#888888")

    # ---- Leaderboard slide ----
    if data.quiz_type == "quiz" and data.leaderboard:
        sl_lb = prs.slides.add_slide(blank_layout)
        add_header_footer(sl_lb, total_slides, total_slides)
        add_text_box(sl_lb, "Leaderboard",
                     Inches(0.3), HEADER_H + Inches(0.05), W - Inches(0.6), Inches(0.6),
                     bold=True, size=22)

        top10 = data.leaderboard[:10]
        if top10:
            cd4 = ChartData()
            cd4.categories = [e.display_name for e in top10]
            cd4.add_series("Score", [e.score for e in top10])
            sl_lb.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED,
                Inches(0.3), HEADER_H + Inches(0.7), W - Inches(0.6), Inches(3.2),
                cd4
            )

        rows = min(len(data.leaderboard), 15) + 1
        tbl = sl_lb.shapes.add_table(rows, 4,
                                     Inches(0.3), HEADER_H + Inches(4.0),
                                     W - Inches(0.6), H - HEADER_H - FOOTER_H - Inches(4.1)).table
        pptx_rank_bg = {1: (0xFF, 0xD7, 0x00), 2: (0xE8, 0xE8, 0xE8), 3: (0xCD, 0x7F, 0x32)}
        for j, hdr in enumerate(["Rank", "Name", "Score", "Time"]):
            cell = tbl.cell(0, j)
            cell.text = hdr
            cell.text_frame.paragraphs[0].runs[0].font.bold = True
            cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x18, 0x90, 0xFF)
        for row_i, e in enumerate(data.leaderboard[:15]):
            tbl.cell(row_i + 1, 0).text = str(e.rank)
            tbl.cell(row_i + 1, 1).text = e.display_name
            tbl.cell(row_i + 1, 2).text = str(e.score)
            time_str = f"{e.time_taken_seconds:.1f}s" if e.time_taken_seconds is not None else "—"
            tbl.cell(row_i + 1, 3).text = time_str
            if e.rank in pptx_rank_bg:
                rgb = pptx_rank_bg[e.rank]
                for j in range(4):
                    c = tbl.cell(row_i + 1, j)
                    c.fill.solid()
                    c.fill.fore_color.rgb = RGBColor(*rgb)
                    c.text_frame.paragraphs[0].runs[0].font.bold = True
            elif row_i % 2 == 0:
                for j in range(4):
                    c = tbl.cell(row_i + 1, j)
                    c.fill.solid()
                    c.fill.fore_color.rgb = RGBColor(0xF5, 0xFA, 0xFF)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# XLSX builder (openpyxl native charts + header/footer)
# ---------------------------------------------------------------------------

def _build_xlsx(data: ExportData) -> bytes:
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment
        from openpyxl.chart import BarChart, PieChart, Reference
        from openpyxl.chart.series import DataPoint
        from openpyxl.drawing.image import Image as XLImage
    except ImportError:
        raise HTTPException(status_code=501, detail="Export libraries not available")

    wb = Workbook()
    header_font = Font(bold=True, color="FFFFFF")
    header_fill = PatternFill("solid", fgColor="1890FF")
    alt_fill_a = PatternFill("solid", fgColor="F5FAFF")
    alt_fill_b = PatternFill("solid", fgColor="FFFFFF")

    def _apply_alt_rows(ws, start_row: int, end_row: int, num_cols: int):
        """Apply alternating row fills from start_row to end_row inclusive."""
        for r in range(start_row, end_row + 1):
            fill = alt_fill_a if r % 2 == 0 else alt_fill_b
            for c in range(1, num_cols + 1):
                ws.cell(r, c).fill = fill

    def _set_col_widths(ws, widths: list):
        """Set column widths by index (1-based)."""
        from openpyxl.utils import get_column_letter
        for i, w in enumerate(widths, 1):
            ws.column_dimensions[get_column_letter(i)].width = w

    logo = _logo_path()

    def _set_sheet_header_footer(ws):
        """Apply print header/footer with page numbers to a worksheet."""
        ws.oddHeader.left.text = "&G"          # logo (embedded separately)
        ws.oddHeader.center.text = "&\"Helvetica,Bold\"&12Swaya.me"
        ws.oddHeader.right.text = data.quiz_title[:40]
        ws.oddFooter.left.text = f"Generated by Swaya.me  |  {data.generated_at.strftime('%d %b %Y')}"
        ws.oddFooter.center.text = "Page &P of &N"
        ws.oddFooter.right.text = f"Session #{data.session_id}"
        ws.sheet_view.showGridLines = True

    # ---- Summary sheet ----
    ws_s = wb.active
    ws_s.title = "Summary"
    _set_sheet_header_footer(ws_s)

    # Logo image in cell A1
    if logo:
        try:
            xl_logo = XLImage(logo)
            xl_logo.width = 80
            xl_logo.height = 40
            ws_s.add_image(xl_logo, "A1")
            ws_s.row_dimensions[1].height = 32
        except Exception:
            pass

    ws_s["C1"] = "Swaya.me  (www.swaya.me)"
    ws_s["C1"].font = Font(bold=True, size=13, color="1890FF")
    ws_s["C1"].hyperlink = 'https://www.swaya.me'
    ws_s["C1"].style = "Hyperlink"

    ws_s["A3"] = "Quiz Title";        ws_s["B3"] = data.quiz_title
    ws_s["A4"] = "Session ID";        ws_s["B4"] = data.session_id
    ws_s["A5"] = "Generated At";      ws_s["B5"] = data.generated_at.strftime("%Y-%m-%d %H:%M")
    ws_s["A6"] = "Total Participants"; ws_s["B6"] = data.total_participants
    ws_s["A7"] = "Total Questions";   ws_s["B7"] = data.total_questions
    ws_s["A8"] = "Quiz Type";         ws_s["B8"] = data.quiz_type.capitalize()

    for row in range(3, 9):
        ws_s[f"A{row}"].font = Font(bold=True, color="1890FF")
        ws_s[f"A{row}"].fill = alt_fill_a if row % 2 == 0 else alt_fill_b
        ws_s[f"B{row}"].fill = alt_fill_a if row % 2 == 0 else alt_fill_b

    _set_col_widths(ws_s, [14, 40, 24])

    if data.quiz_type == "quiz":
        total_correct = sum(
            q.answer_distribution[q.correct_answer_index]
            for q in data.questions
            if q.options and q.correct_answer_index is not None
               and q.correct_answer_index < len(q.answer_distribution)
        )
        total_answers = sum(sum(q.answer_distribution) for q in data.questions if q.options)
        total_wrong = total_answers - total_correct

        ws_s["A10"] = "Correct Answers"; ws_s["B10"] = total_correct
        ws_s["A11"] = "Wrong Answers";   ws_s["B11"] = total_wrong
        ws_s["A10"].font = Font(bold=True)
        ws_s["A11"].font = Font(bold=True)

        if total_correct + total_wrong > 0:
            pc = PieChart()
            pc.title = "Correct vs Wrong"
            labels = Reference(ws_s, min_col=1, min_row=10, max_row=11)
            data_ref = Reference(ws_s, min_col=2, min_row=10, max_row=11)
            pc.add_data(data_ref)
            pc.set_categories(labels)
            pc.width = 10; pc.height = 8
            ws_s.add_chart(pc, "D4")

    # ---- Questions sheet ----
    ws_q = wb.create_sheet("Questions")
    _set_sheet_header_footer(ws_q)
    q_headers = ["Q#", "Text", "Type", "Opt A", "Opt B", "Opt C", "Opt D",
                 "Votes A", "Votes B", "Votes C", "Votes D",
                 "% A", "% B", "% C", "% D", "Correct Option"]
    for col, h in enumerate(q_headers, 1):
        cell = ws_q.cell(row=1, column=col, value=h)
        cell.font = header_font
        cell.fill = header_fill

    for idx, q in enumerate(data.questions):
        r = idx + 2
        total = sum(q.answer_distribution) or 1
        ws_q.cell(r, 1, idx + 1)
        ws_q.cell(r, 2, q.text)
        ws_q.cell(r, 3, q.question_type)
        if q.options:
            for i, opt in enumerate(q.options[:4]):
                ws_q.cell(r, 4 + i, opt)
            for i, cnt in enumerate(q.answer_distribution[:4]):
                ws_q.cell(r, 8 + i, cnt)
                ws_q.cell(r, 12 + i, round(cnt / total * 100, 1))
            if q.correct_answer_index is not None:
                ws_q.cell(r, 16, chr(65 + q.correct_answer_index))

        # Question image
        if q.question_image_path:
            try:
                xl_qimg = XLImage(q.question_image_path)
                xl_qimg.width = 120
                xl_qimg.height = 80
                # Place image in column P of this row
                ws_q.add_image(xl_qimg, f"P{r}")
                ws_q.row_dimensions[r].height = 65
            except Exception:
                pass
        # Option images
        if q.option_image_paths:
            for i, opt_key in enumerate(("A", "B", "C", "D")):
                img_path = q.option_image_paths.get(opt_key)
                if img_path:
                    try:
                        xl_oimg = XLImage(img_path)
                        xl_oimg.width = 80
                        xl_oimg.height = 60
                        col_letter = chr(ord("Q") + i)
                        ws_q.add_image(xl_oimg, f"{col_letter}{r}")
                        ws_q.row_dimensions[r].height = max(ws_q.row_dimensions[r].height or 0, 50)
                    except Exception:
                        pass

    # Alternating rows + col widths for Questions sheet
    if data.questions:
        _apply_alt_rows(ws_q, 2, len(data.questions) + 1, len(q_headers))
    _set_col_widths(ws_q, [4, 40, 12, 20, 20, 20, 20, 8, 8, 8, 8, 7, 7, 7, 7, 12])

    if len(data.questions) > 0:
        ws_q.cell(1, 18, "Q")
        ws_q.cell(1, 19, "Total Votes")
        ws_q.cell(1, 18).font = header_font
        ws_q.cell(1, 18).fill = header_fill
        ws_q.cell(1, 19).font = header_font
        ws_q.cell(1, 19).fill = header_fill
        for idx, q in enumerate(data.questions):
            ws_q.cell(idx + 2, 18, f"Q{idx+1}")
            ws_q.cell(idx + 2, 19, sum(q.answer_distribution))

        bc = BarChart()
        bc.type = "bar"
        bc.title = "Responses per Question"
        bc.y_axis.title = "Votes"
        end_row = len(data.questions) + 1
        data_ref = Reference(ws_q, min_col=19, min_row=1, max_row=end_row)
        cats_ref = Reference(ws_q, min_col=18, min_row=2, max_row=end_row)
        bc.add_data(data_ref, titles_from_data=True)
        bc.set_categories(cats_ref)
        bc.width = 16; bc.height = 12
        # Apply brand color to bar series
        try:
            from openpyxl.chart.data_source import NumDataSource
            ser = bc.series[0]
            ser.graphicalProperties.solidFill = "1890FF"
        except Exception:
            pass
        ws_q.add_chart(bc, "R2")

    # ---- Leaderboard sheet ----
    if data.quiz_type == "quiz":
        ws_lb = wb.create_sheet("Leaderboard")
        _set_sheet_header_footer(ws_lb)
        lb_headers = ["Rank", "Name", "Score", "Time (s)"]
        for col, h in enumerate(lb_headers, 1):
            cell = ws_lb.cell(row=1, column=col, value=h)
            cell.font = header_font
            cell.fill = header_fill

        rank_fills = {
            1: PatternFill("solid", fgColor="FFD700"),
            2: PatternFill("solid", fgColor="E8E8E8"),
            3: PatternFill("solid", fgColor="CD7F32"),
        }
        rank_fonts = {
            1: Font(bold=True, color="7A5700"),
            2: Font(bold=True, color="555555"),
            3: Font(bold=True, color="5C2E00"),
        }
        for i, e in enumerate(data.leaderboard):
            r = i + 2
            ws_lb.cell(r, 1, e.rank)
            ws_lb.cell(r, 2, e.display_name)
            ws_lb.cell(r, 3, e.score)
            ws_lb.cell(r, 4, e.time_taken_seconds if e.time_taken_seconds is not None else "")
            if e.rank in rank_fills:
                for col in range(1, 5):
                    ws_lb.cell(r, col).fill = rank_fills[e.rank]
                    ws_lb.cell(r, col).font = rank_fonts[e.rank]
            else:
                fill = alt_fill_a if i % 2 == 0 else alt_fill_b
                for col in range(1, 5):
                    ws_lb.cell(r, col).fill = fill

        _set_col_widths(ws_lb, [6, 30, 10, 12])

        if data.leaderboard:
            bc2 = BarChart()
            bc2.type = "bar"
            bc2.title = "Top Participant Scores"
            end_lb = min(len(data.leaderboard), 10) + 1
            data_ref2 = Reference(ws_lb, min_col=3, min_row=1, max_row=end_lb)
            cats_ref2 = Reference(ws_lb, min_col=2, min_row=2, max_row=end_lb)
            bc2.add_data(data_ref2, titles_from_data=True)
            bc2.set_categories(cats_ref2)
            bc2.width = 16; bc2.height = 12
            try:
                ser2 = bc2.series[0]
                ser2.graphicalProperties.solidFill = "52C41A"
            except Exception:
                pass
            ws_lb.add_chart(bc2, "F2")

    # ---- Word Cloud sheet ----
    wc_questions = [q for q in data.questions if q.word_frequencies]
    if wc_questions:
        ws_wc = wb.create_sheet("Word Cloud Data")
        _set_sheet_header_footer(ws_wc)
        for col, h in enumerate(["Question", "Word", "Count", "%"], 1):
            cell = ws_wc.cell(row=1, column=col, value=h)
            cell.font = header_font
            cell.fill = header_fill

        row = 2
        img_anchor_row = 2
        for q in wc_questions:
            total_wc = sum(q.word_frequencies.values()) or 1
            for word, cnt in sorted(q.word_frequencies.items(), key=lambda x: -x[1]):
                ws_wc.cell(row, 1, q.text[:60])
                ws_wc.cell(row, 2, word)
                ws_wc.cell(row, 3, cnt)
                ws_wc.cell(row, 4, round(cnt / total_wc * 100, 1))
                row += 1

            # Embed word cloud image to the right of the data
            wc_png = _make_wordcloud_png(q.word_frequencies, width=500, height=250)
            if wc_png:
                try:
                    xl_wc = XLImage(io.BytesIO(wc_png))
                    xl_wc.width = 400
                    xl_wc.height = 200
                    ws_wc.add_image(xl_wc, f"F{img_anchor_row}")
                    img_anchor_row += 18
                except Exception:
                    pass

    # ---- Raw Data sheet ----
    ws_raw = wb.create_sheet("Raw Data")
    _set_sheet_header_footer(ws_raw)
    raw_headers = ["Q#", "Question Text", "Type", "Option / Response", "Option Index", "Votes", "% of Total"]
    for col, h in enumerate(raw_headers, 1):
        cell = ws_raw.cell(row=1, column=col, value=h)
        cell.font = header_font
        cell.fill = header_fill

    raw_row = 2
    for idx, q in enumerate(data.questions):
        if q.options:
            # MCQ / scale with labelled options
            total = sum(q.answer_distribution) or 1
            for i, (opt, cnt) in enumerate(zip(q.options, q.answer_distribution)):
                ws_raw.cell(raw_row, 1, idx + 1)
                ws_raw.cell(raw_row, 2, q.text)
                ws_raw.cell(raw_row, 3, q.question_type)
                ws_raw.cell(raw_row, 4, opt)
                ws_raw.cell(raw_row, 5, i)
                ws_raw.cell(raw_row, 6, cnt)
                ws_raw.cell(raw_row, 7, round(cnt / total * 100, 1))
                raw_row += 1
        elif q.text_answers:
            # single_line / paragraph — one row per response
            for answer_text in q.text_answers:
                ws_raw.cell(raw_row, 1, idx + 1)
                ws_raw.cell(raw_row, 2, q.text)
                ws_raw.cell(raw_row, 3, q.question_type)
                ws_raw.cell(raw_row, 4, answer_text)
                ws_raw.cell(raw_row, 5, None)
                ws_raw.cell(raw_row, 6, None)
                ws_raw.cell(raw_row, 7, None)
                raw_row += 1

    # Alternating rows + col widths for Raw Data sheet
    if raw_row > 2:
        _apply_alt_rows(ws_raw, 2, raw_row - 1, len(raw_headers))
    _set_col_widths(ws_raw, [4, 45, 12, 35, 10, 8, 8])

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Main ExportService
# ---------------------------------------------------------------------------

class ExportService:

    async def generate(
        self,
        session_id: int,
        fmt: str,
        db: AsyncSession,
        tenant_id: int,
        answer_service: AnswerServiceAsync,
    ) -> Tuple[bytes, str, str]:
        """
        Generate export file.

        Returns (file_bytes, media_type, filename).
        """
        try:
            data = await self._gather_data(session_id, db, tenant_id, answer_service)
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Failed to gather export data: {e}")

        dt_str = data.generated_at.strftime("%Y%m%d_%H%M%S")
        safe_title = "".join(c if c.isalnum() or c in "_-" else "_" for c in data.quiz_title).strip("_")[:40]
        base_name = f"SwayameExtract_{safe_title}_{dt_str}"

        try:
            if fmt == "pdf":
                file_bytes = _build_pdf(data)
                media_type = "application/pdf"
                filename = f"{base_name}.pdf"
            elif fmt == "docx":
                file_bytes = _build_docx(data)
                media_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                filename = f"{base_name}.docx"
            elif fmt == "pptx":
                file_bytes = _build_pptx(data)
                media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                filename = f"{base_name}.pptx"
            elif fmt == "xlsx":
                file_bytes = _build_xlsx(data)
                media_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                filename = f"{base_name}.xlsx"
            else:
                raise HTTPException(status_code=400, detail="Unsupported format")
        except HTTPException:
            raise
        except ImportError as e:
            raise HTTPException(status_code=501, detail="Export libraries not available")
        except Exception as e:
            import traceback as _tb
            import logging as _log
            _log.getLogger(__name__).error("Export failed (fmt=%s session=%s): %s", fmt, session_id, _tb.format_exc())
            raise HTTPException(status_code=500, detail=f"Failed to generate export: {e}")

        return file_bytes, media_type, filename

    async def _gather_data(
        self,
        session_id: int,
        db: AsyncSession,
        tenant_id: int,
        answer_service: AnswerServiceAsync,
    ) -> ExportData:
        result = await db.execute(
            select(QuizSession)
            .filter(QuizSession.id == session_id)
            .options(joinedload(QuizSession.quiz).selectinload(Quiz.questions))
        )
        session = result.scalar_one_or_none()

        if not session:
            raise HTTPException(status_code=404, detail="Session not found")
        if session.quiz.tenant_id != tenant_id:
            raise HTTPException(status_code=403, detail="Access denied")

        quiz = session.quiz
        questions = sorted(quiz.questions, key=lambda q: q.order)

        session_results = await answer_service.get_session_results(db, session_id)
        results_by_qid = {qr.question_id: qr for qr in session_results.question_results}

        lb_response = await answer_service.get_leaderboard(db, session_id)
        leaderboard = [
            LeaderboardEntryExport(
                rank=e.rank,
                display_name=e.display_name,
                score=e.score,
                time_taken_seconds=e.time_taken_seconds,
            )
            for e in lb_response.entries
        ]

        q_exports = []
        for q in questions:
            qr = results_by_qid.get(q.id)
            dist = qr.answer_distribution if qr else []
            total_ans = qr.total_answers if qr else 0

            wf = None
            if q.question_type == QuestionType.WORD_CLOUD:
                try:
                    wc = await answer_service.get_word_cloud_results(db, session_id, q.id)
                    wf = wc.word_frequencies
                except Exception:
                    wf = {}

            ta = None
            if q.question_type in (QuestionType.SINGLE_LINE, QuestionType.PARAGRAPH):
                try:
                    rows = await db.execute(
                        select(Answer.text_answer)
                        .filter(
                            Answer.session_id == session_id,
                            Answer.question_id == q.id,
                            Answer.text_answer.isnot(None),
                        )
                        .order_by(Answer.created_at)
                    )
                    ta = [r for (r,) in rows.all() if r and r.strip()]
                except Exception:
                    ta = []

            q_img = _resolve_upload_path(q.question_image_url)
            opt_imgs: Dict[str, str] = {}
            if q.option_images:
                for key, path in q.option_images.items():
                    rp = _resolve_upload_path(path)
                    if rp:
                        opt_imgs[key] = rp

            q_exports.append(QuestionExport(
                id=q.id,
                text=q.text,
                question_type=q.question_type.value,
                points=q.points or 1,
                max_time_seconds=q.max_time_seconds,
                options=q.options if q.options else None,
                correct_answer_index=q.correct_answer_index,
                answer_distribution=dist,
                total_answers=total_ans,
                word_frequencies=wf,
                text_answers=ta,
                question_image_path=q_img,
                option_image_paths=opt_imgs if opt_imgs else None,
            ))

        return ExportData(
            session_id=session_id,
            quiz_title=quiz.title,
            quiz_type=quiz.quiz_type.value,
            total_questions=len(questions),
            total_participants=session_results.total_participants,
            generated_at=datetime.utcnow(),
            questions=q_exports,
            leaderboard=leaderboard,
        )
